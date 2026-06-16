#!/usr/bin/env python3
"""Print a school template's layouts, placeholders, and theme fonts.

Use the layout names this prints when you fill in lesson.json, so the generated deck
lands on the school's real layouts instead of a generic guess.

Usage:
    python3 inspect_template.py school.pptx
"""
import sys

from pptx import Presentation


def main():
    if len(sys.argv) != 2:
        sys.exit("usage: inspect_template.py <template.pptx>")
    prs = Presentation(sys.argv[1])

    print(f"Slide size: {prs.slide_width} x {prs.slide_height} EMU")
    print(f"Layouts ({len(prs.slide_layouts)}):")
    for i, layout in enumerate(prs.slide_layouts):
        phs = []
        for ph in layout.placeholders:
            try:
                ph_type = ph.placeholder_format.type
            except ValueError:
                ph_type = "?"
            phs.append(f"{ph.placeholder_format.idx}:{ph_type}")
        ph_str = ", ".join(phs) if phs else "(no placeholders)"
        print(f"  [{i}] {layout.name}")
        print(f"       placeholders -> {ph_str}")

    # Theme fonts from the first master, if present.
    try:
        master = prs.slide_masters[0]
        font_scheme = master.element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}fontScheme"
        )
        if font_scheme is not None:
            print(f"Theme font scheme: {font_scheme.get('name')}")
    except Exception:  # noqa: BLE001
        pass


if __name__ == "__main__":
    main()
