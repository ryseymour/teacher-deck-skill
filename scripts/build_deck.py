#!/usr/bin/env python3
"""Build a slide deck from a structured lesson, reusing a school's .pptx template.

The template supplies the theme, fonts, colors, and slide layouts. This script clears
the template's example slides (keeping its design) and adds new slides populated from a
lesson JSON. Teacher notes and image attribution are written into each slide's notes pane.

Usage:
    python3 build_deck.py --template school.pptx --lesson lesson.json --out deck.pptx
    python3 build_deck.py --template school.pptx --lesson lesson.json --out deck.pptx --keep-template-slides

See references/lesson_schema.md for the lesson JSON shape.
"""
import argparse
import json
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor


def load_lesson(path):
    with open(path, "r", encoding="utf-8") as fh:
        return json.load(fh)


def clear_slides(prs):
    """Remove existing slides but keep masters, layouts, and theme."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        sld_id_lst.remove(sld_id)


def pick_layout(prs, ref, default_idx=1):
    """Resolve a layout by name (case-insensitive) or integer index, with a fallback."""
    layouts = list(prs.slide_layouts)
    if ref is None:
        return layouts[min(default_idx, len(layouts) - 1)]
    if isinstance(ref, int):
        if 0 <= ref < len(layouts):
            return layouts[ref]
        return layouts[min(default_idx, len(layouts) - 1)]
    for layout in layouts:
        if layout.name.strip().lower() == str(ref).strip().lower():
            return layout
    # Loose contains-match before giving up
    for layout in layouts:
        if str(ref).strip().lower() in layout.name.strip().lower():
            return layout
    sys.stderr.write(
        f"  ! layout '{ref}' not found; using '{layouts[min(default_idx, len(layouts) - 1)].name}'\n"
    )
    return layouts[min(default_idx, len(layouts) - 1)]


def set_title(slide, title):
    if title and slide.shapes.title is not None:
        slide.shapes.title.text = title


def find_body_placeholder(slide):
    """Return the best placeholder for bullet content, skipping the title."""
    title = slide.shapes.title
    # Prefer BODY/OBJECT/CONTENT placeholders.
    preferred = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
    candidates = []
    for ph in slide.placeholders:
        if title is not None and ph._element is title._element:
            continue
        if not ph.has_text_frame:
            continue
        candidates.append(ph)
    for ph in candidates:
        try:
            if ph.placeholder_format.type in preferred:
                return ph
        except ValueError:
            continue
    return candidates[0] if candidates else None


def add_bullets(slide, bullets):
    if not bullets:
        return
    body = find_body_placeholder(slide)
    if body is None:
        return
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = str(line)


def add_image(slide, image_path):
    """Insert into a picture placeholder if present, else float a sized picture."""
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                ph.insert_picture(image_path)
                return
        except (ValueError, Exception):
            continue
    # No picture placeholder: place a reasonably sized image on the lower right.
    try:
        slide.shapes.add_picture(image_path, Inches(5.2), Inches(2.2), height=Inches(3.6))
    except Exception as exc:  # noqa: BLE001
        sys.stderr.write(f"  ! could not add image {image_path}: {exc}\n")


def write_notes(slide, parts):
    text = "\n".join(p for p in parts if p)
    if not text:
        return
    slide.notes_slide.notes_text_frame.text = text


def requires_visible_credit(credit):
    """True for licenses that require visible attribution (CC-BY, CC-BY-SA, etc.)."""
    if not credit:
        return False
    c = credit.lower()
    if "cc0" in c or "public domain" in c or "pdm" in c:
        return False
    return "cc by" in c or "cc-by" in c


def add_credit_caption(slide, prs, text):
    """Place a small, gray attribution line along the bottom of the slide."""
    left = Inches(0.3)
    width = prs.slide_width - Inches(0.6)
    top = prs.slide_height - Inches(0.45)
    box = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.text = text
    run = para.runs[0]
    run.font.size = Pt(8)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)


def build(template, lesson, out, keep_template_slides=False, credit_on_slide="auto"):
    prs = Presentation(template)
    if not keep_template_slides:
        clear_slides(prs)

    slides = lesson.get("slides", [])
    if not slides:
        sys.exit("lesson JSON has no 'slides'.")

    for i, spec in enumerate(slides, 1):
        layout = pick_layout(prs, spec.get("layout"))
        slide = prs.slides.add_slide(layout)
        set_title(slide, spec.get("title"))
        add_bullets(slide, spec.get("bullets"))

        image = spec.get("image")
        if image:
            add_image(slide, image)

        notes_parts = []
        if spec.get("notes"):
            notes_parts.append(spec["notes"])
        if spec.get("check_for_understanding"):
            notes_parts.append("Check for understanding: " + spec["check_for_understanding"])
        credit = spec.get("image_credit")
        if credit:
            notes_parts.append("Image: " + credit)
        write_notes(slide, notes_parts)

        if image and credit:
            show = credit_on_slide == "always" or (
                credit_on_slide == "auto" and requires_visible_credit(credit)
            )
            if show:
                add_credit_caption(slide, prs, credit)

        print(f"  slide {i}: '{spec.get('title', '(untitled)')}' on layout '{layout.name}'")

    prs.save(out)
    print(f"Saved {len(slides)} slides to {out}")


def main():
    ap = argparse.ArgumentParser(description="Build a school-templated deck from a lesson JSON.")
    ap.add_argument("--template", required=True, help="Path to the school's example .pptx")
    ap.add_argument("--lesson", required=True, help="Path to lesson.json")
    ap.add_argument("--out", required=True, help="Path for the output .pptx")
    ap.add_argument(
        "--keep-template-slides",
        action="store_true",
        help="Keep the template's existing slides instead of replacing them.",
    )
    ap.add_argument(
        "--credit-on-slide",
        choices=["auto", "always", "off"],
        default="auto",
        help="Show image attribution on the slide face. 'auto' (default) does it only for "
        "licenses that require visible credit (CC-BY); 'always' for every credited image; "
        "'off' keeps all credit in the notes only.",
    )
    args = ap.parse_args()
    build(
        args.template,
        load_lesson(args.lesson),
        args.out,
        args.keep_template_slides,
        args.credit_on_slide,
    )


if __name__ == "__main__":
    main()
